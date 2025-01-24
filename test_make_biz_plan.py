import unittest
import json
from pptx.util import Inches
from unittest.mock import patch, MagicMock
from Make_Biz_Plan import OKRId, add_goal_image, flip_bool_attribute, SquareDimensions, LineDimensions, VivaGoal, get_goal_by_id, get_parent_goals_from_alignment, get_theme_goal_by_id, create_slide, add_goal_details_to_slide, add_text_block_to_slide, ACTION_TYPE, OUTCOME_TYPE, goal_sort_key, load_goals_from_workbook

class TestUtilityFunctions(unittest.TestCase):
    def test_flip_bool_attribute(self):
        class DummyObject:
            def __init__(self):
                self.test_attr = True
        obj = DummyObject()
        original_value = obj.test_attr
        flip_bool_attribute(obj, 'test_attr')
        self.assertEqual(obj.test_attr, original_value)

    def test_square_dimensions(self):
        dims = SquareDimensions(1, 2, 3, 4)
        self.assertEqual(dims.left, Inches(1))
        self.assertEqual(dims.top, Inches(2))
        self.assertEqual(dims.width, Inches(3))
        self.assertEqual(dims.height, Inches(4))

    def test_line_dimensions(self):
        dims = LineDimensions(1, 2, 3)
        self.assertEqual(dims.left, Inches(1))
        self.assertEqual(dims.top, Inches(2))
        self.assertEqual(dims.width, Inches(3))

class TestVivaGoal(unittest.TestCase):
    def setUp(self):
        self.headers = ['Id', 'Title', 'Tag', 'Owner', 'Period', 'Start Date', 'End Date',
                        'Description', 'Aligned To (weight, Objective ID)', 'Metric Name',
                        'Target', 'Object Type', 'Status']

        # Base template for a row
        base_row = ['', '', '', 'John Doe', 'Q1 2024', '2024-01-01', '2024-03-31',
                    'Description text', '', 'Metric', 'Target', '', 'On Track']

        # Create a specific test row for basic creation test
        self.test_row = base_row.copy()
        self.test_row[0:3] = ['1', 'Test Goal', 'Theme']
        self.test_row[3] = 'John'  # Specific owner for the test
        self.test_row[11] = 'Objective'

        self.rows = []

        # Theme 1 hierarchy
        theme1 = base_row.copy()
        theme1[0:3] = ['1', 'Theme 1', 'Theme']
        theme1[11] = 'Objective'

        outcome1 = base_row.copy()
        outcome1[0:3] = ['2', 'Outcome 1', '']
        outcome1[8] = '(weight: 100%, Id: 1)'
        outcome1[11] = 'Outcome'

        objective1 = base_row.copy()
        objective1[0:3] = ['3', 'Objective 1', '']
        objective1[8] = '(weight: 100%, Id: 1)'
        objective1[11] = 'Objective'

        action1 = base_row.copy()
        action1[0:3] = ['4', 'Action 1', '']
        action1[8] = '(weight: 100%, Id: 3)'
        action1[11] = 'Action'

        # Theme 2 hierarchy with multiple outcomes and actions
        theme2 = base_row.copy()
        theme2[0:3] = ['5', 'Theme 2', 'Theme']
        theme2[11] = 'Objective'

        objective2 = base_row.copy()
        objective2[0:3] = ['6', 'Objective 2', '']
        objective2[8] = '(weight: 100%, Id: 5)'
        objective2[11] = 'Objective'

        outcome2 = base_row.copy()
        outcome2[0:3] = ['7', 'Outcome 2', '']
        outcome2[8] = '(weight: 100%, Id: 6)'
        outcome2[11] = 'Outcome'

        action2 = base_row.copy()
        action2[0:3] = ['8', 'Action 2', '']
        action2[8] = '(weight: 100%, Id: 6)'
        action2[11] = 'Action'

        # Complex hierarchy with multiple levels
        theme3 = base_row.copy()
        theme3[0:3] = ['9', 'Theme 3', 'Theme']
        theme3[11] = 'Objective'

        outcome3a = base_row.copy()
        outcome3a[0:3] = ['10', 'Outcome 3A', '']
        outcome3a[8] = '(weight: 100%, Id: 9)'
        outcome3a[11] = 'Outcome'

        objective3a = base_row.copy()
        objective3a[0:3] = ['11', 'Objective 3A', '']
        objective3a[8] = '(weight: 100%, Id: 9)'
        objective3a[11] = 'Objective'

        outcome3b = base_row.copy()
        outcome3b[0:3] = ['12', 'Outcome 3B', '']
        outcome3b[8] = '(weight: 100%, Id: 11)'
        outcome3b[11] = 'Outcome'

        action3a = base_row.copy()
        action3a[0:3] = ['13', 'Action 3A', '']
        action3a[8] = '(weight: 100%, Id: 11)'
        action3a[11] = 'Action'

        # Add rows in an order that will test the sorting
        self.rows = [
            action3a,     # Should appear last in its hierarchy
            theme1,       # Should appear first
            outcome2,     # Should appear before action2
            objective3a,  # Should appear after outcome3a
            theme2,       # Should appear after theme1's hierarchy
            outcome3b,    # Should appear before action3a
            action1,      # Should appear after objective1
            objective2,   # Should appear after theme2
            theme3,       # Should appear last
            outcome1,     # Should appear before objective1
            objective1,   # Should appear after outcome1
            action2,      # Should appear after outcome2
            outcome3a,    # Should appear before objective3a
        ]

        # The row used for basic tests
        self.row = self.test_row

    def test_viva_goal_creation(self):
        goal = VivaGoal(self.row, self.headers, 0)
        self.assertEqual(goal.okr_id, '1')
        self.assertEqual(goal.title, 'Test Goal')
        self.assertEqual(goal.tag, 'Theme')
        self.assertEqual(goal.owner, 'John')

    def test_viva_goal_invalid_row(self):
        # Test that IndexError is raised when row is too short
        with self.assertRaises(IndexError):
            VivaGoal(['1'], self.headers, 0)

class TestOKRId(unittest.TestCase):
    def test_valid_okr_id(self):
        okr = OKRId('"http://example.com" "123"')
        self.assertEqual(okr.okr_link, 'http://example.com')
        self.assertEqual(okr.okr_id, '123')

    def test_invalid_okr_id(self):
        okr = OKRId('invalid')
        self.assertEqual(okr.okr_link, '')
        self.assertEqual(okr.okr_id, '')

class TestGoalOperations(unittest.TestCase):
    def setUp(self):
        # Store the original goals_dict
        import Make_Biz_Plan
        self.original_goals_dict = Make_Biz_Plan.goals_dict
        # Replace with test data
        Make_Biz_Plan.goals_dict = {'1': 'test_goal'}

    def tearDown(self):
        # Restore the original goals_dict
        import Make_Biz_Plan
        Make_Biz_Plan.goals_dict = self.original_goals_dict

    def test_get_goal_by_id(self):
        self.assertEqual(get_goal_by_id('1'), 'test_goal')
        self.assertIsNone(get_goal_by_id('nonexistent'))

    def test_get_theme_goal_by_id(self):
        goals = [
            MagicMock(tag='Other'),
            MagicMock(tag='Theme'),
            MagicMock(tag='Other')
        ]
        theme_goal = get_theme_goal_by_id(goals)
        self.assertEqual(theme_goal.tag, 'Theme')

class TestPresentationOperations(unittest.TestCase):
    @patch('pptx.Presentation')
    def test_create_slide(self, mock_presentation):
        slide = create_slide(mock_presentation, (0, 0), 'Test Title')
        mock_presentation.slide_masters.__getitem__.assert_called_once()
        self.assertIsNotNone(slide)

    @patch('Make_Biz_Plan.add_text_block_to_slide')
    def test_add_goal_details_to_slide(self, mock_add_text):
        slide = MagicMock()
        goal = MagicMock(
            object_type='Objective',
            metric_name='Test Metric',
            target='100%',
            owner='John Doe',
            schedule='Q1',
            status='On Track'
        )
        add_goal_details_to_slide(slide, goal)
        mock_add_text.assert_called_once()

class TestSorting(unittest.TestCase):
    def setUp(self):
        # Initialize TestVivaGoal to get test data
        self.test_viva_goal = TestVivaGoal()
        self.test_viva_goal.setUp()

        # Create the goals and populate goals_dict
        import Make_Biz_Plan
        self.original_goals_dict = Make_Biz_Plan.goals_dict
        Make_Biz_Plan.goals_dict = {}  # Reset goals_dict

        self.goals = []
        for idx, row in enumerate(self.test_viva_goal.rows):
            goal = VivaGoal(row, self.test_viva_goal.headers, idx)
            self.goals.append(goal)
            # Add to goals_dict using the ID from the row
            Make_Biz_Plan.goals_dict[row[0]] = goal

    def tearDown(self):
        # Restore original goals_dict
        import Make_Biz_Plan
        Make_Biz_Plan.goals_dict = self.original_goals_dict

    def test_goal_sorting(self):
        """Test that goals are sorted in the correct order"""
        sorted_goals = sorted(self.goals, key=goal_sort_key)

        # Verify the order follows the expected pattern:
        # Theme -> [Outcome] -> Objective -> [Outcome] -> Action
        expected_order = [
            'Theme 1',
            'Outcome 1',
            'Objective 1',
            'Action 1',
            'Theme 2',
            'Objective 2',
            'Outcome 2',
            'Action 2',
            'Theme 3',
            'Outcome 3A',
            'Objective 3A',
            'Outcome 3B',
            'Action 3A'
        ]

        actual_order = [goal.title for goal in sorted_goals]
        self.assertEqual(actual_order, expected_order)

class TestExceptionHandling(unittest.TestCase):
    def setUp(self):
        # Initialize test data
        self.headers = ['Id', 'Title', 'Tag', 'Owner', 'Period', 'Start Date', 'End Date',
                        'Description', 'Aligned To (weight, Objective ID)', 'Metric Name',
                        'Target', 'Object Type', 'Status']

        # Import and store original goals_dict
        import Make_Biz_Plan
        self.Make_Biz_Plan = Make_Biz_Plan  # Store reference to module
        self.original_goals_dict = Make_Biz_Plan.goals_dict
        Make_Biz_Plan.goals_dict = {}

    def tearDown(self):
        # Restore original goals_dict
        self.Make_Biz_Plan.goals_dict = self.original_goals_dict

    def test_add_text_block_invalid_json(self):
        """Test handling of invalid JSON in add_text_block_to_slide"""
        text_frame = MagicMock()
        with self.assertRaises(json.JSONDecodeError):
            add_text_block_to_slide(text_frame, "invalid json")

    def test_add_text_block_missing_elements(self):
        """Test handling of JSON without elements array"""
        text_frame = MagicMock()
        with self.assertRaises(KeyError):
            add_text_block_to_slide(text_frame, '{"wrong_key": []}')

    def test_create_slide_invalid_master(self):
        """Test handling of invalid slide master index"""
        prs = MagicMock()
        prs.slide_masters.__getitem__.side_effect = IndexError
        with self.assertRaises(ValueError):
            create_slide(prs, (999, 0), "Test Title")

    def test_goal_sort_key_invalid_parent_reference(self):
        """Test handling of invalid parent references in goal sorting"""
        goal = MagicMock(
            object_type=ACTION_TYPE,
            alignment='(weight: 100%, Id: nonexistent)',
            title='Test Action'
        )
        with self.assertRaises(ValueError):
            goal_sort_key(goal)

    def test_multiple_parent_goals_error(self):
        """Test handling of multiple parent goals for outcomes/actions"""
        goal = MagicMock(
            object_type=OUTCOME_TYPE,
            alignment='(weight: 50%, Id: 1)(weight: 50%, Id: 2)',
            title='Test Outcome'
        )
        self.Make_Biz_Plan.goals_dict = {'1': MagicMock(), '2': MagicMock()}
        with self.assertRaises(ValueError):
            goal_sort_key(goal)

    @patch('openpyxl.load_workbook')
    def test_load_goals_invalid_workbook(self, mock_load):
        """Test handling of invalid Excel workbook"""
        mock_load.side_effect = Exception("Invalid workbook")
        with self.assertRaises(ValueError):
            load_goals_from_workbook("invalid.xlsx")

    def test_add_goal_image_missing_file(self):
        """Test handling of missing image file"""
        slide = MagicMock()
        goal = MagicMock()
        with self.assertRaises(ValueError):
            add_goal_image(slide, goal, "nonexistent.png")

class TestVivaGoalValidation(unittest.TestCase):
    def setUp(self):
        self.headers = ['Id', 'Title', 'Tag', 'Owner', 'Period', 'Start Date', 'End Date',
                        'Description', 'Aligned To (weight, Objective ID)', 'Metric Name',
                        'Target', 'Object Type', 'Status']
        self.base_row = ['1', 'Test Goal', 'Theme', 'John', 'Q1', '2024-01-01', '2024-03-31',
                         'Description', '', 'Metric', 'Target', 'Objective', 'On Track']

    def test_invalid_object_type(self):
        """Test handling of invalid object type"""
        row = self.base_row.copy()
        row[self.headers.index('Object Type')] = 'InvalidType'
        goal = VivaGoal(row, self.headers, 0)
        with self.assertRaises(ValueError):
            goal_sort_key(goal)

    def test_invalid_alignment_format(self):
        """Test handling of malformed alignment string"""
        row = self.base_row.copy()
        row[self.headers.index('Aligned To (weight, Objective ID)')] = 'malformed alignment'
        goal = VivaGoal(row, self.headers, 0)
        self.assertEqual(get_parent_goals_from_alignment(goal), [])

    def test_invalid_date_format(self):
        """Test handling of invalid date format"""
        row = self.base_row.copy()
        row[self.headers.index('Start Date')] = 'not a date'
        row[self.headers.index('End Date')] = 'also not a date'
        goal = VivaGoal(row, self.headers, 0)
        # Verify that the goal is created but dates are stored as strings
        self.assertEqual(goal.start_date, 'not a date')
        self.assertEqual(goal.end_date, 'also not a date')


if __name__ == '__main__':
    unittest.main()
