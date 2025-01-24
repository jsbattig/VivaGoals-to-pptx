# This scripts transforms an export from Viva Goals into a PowerPoint file with one slide per Viva Goals object

from pptx import Presentation
from openpyxl import load_workbook
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import re
import json
import argparse

SOURCE_WORKBOOK = 'VivaGoals.xlsx'
TEMPLATE_POWERPOINT = 'template.pptx'
OBJECTIVE_IMAGE = 'objective.png'
INITIATIVE_IMAGE = 'initiative.png'
OUTCOME_IMAGE = 'outcome.png'
TARGET_BIZPLAN_POWERPOINT = 'bizplan.pptx'

THEME_TAG = "Theme"

OBJECTIVE_TYPE = "Objective"
OUTCOME_TYPE = "Outcome"
ACTION_TYPE = "Action"

THEME_SLIDE_MASTER = 0
THEME_SLIDE_MASTER_LAYOUT = 3
OKR_SLIDE_MASTER = 2
OKR_SLIDE_MASTER_LAYOUT = 11

class SquareDimensions:
    def __init__(self, left, top, width, height):
        self.left = Inches(left)
        self.top = Inches(top)
        self.width = Inches(width)
        self.height = Inches(height)
        
class LineDimensions:
    def __init__(self, left, top, width):
        self.left = Inches(left)
        self.top = Inches(top)
        self.width = Inches(width)

# Iterate over each row in the Excel file starting from the second row (first row has headers)
class VivaGoal:
    def __init__(self, row, headers, row_number):
        self.okr_id = row[headers.index('Id')]
        self.title = row[headers.index('Title')]
        self.tag = row[headers.index('Tag')]
        self.owner = row[headers.index('Owner')]
        self.schedule = row[headers.index('Period')]
        self.start_date = row[headers.index('Start Date')]
        self.end_date = row[headers.index('End Date')]
        self.description = row[headers.index('Description')]
        self.alignment = row[headers.index('Aligned To (weight, Objective ID)')]
        self.metric_name = row[headers.index('Metric Name')]
        self.target = row[headers.index('Target')]
        self.object_type = row[headers.index('Object Type')]
        self.status = row[headers.index('Status')]
        self.row_number = row_number  # Add row number attribute

class OKRId:
    def __init__(self, okr_id_str):
        matches = re.findall(r'"(.*?)"', okr_id_str)
        if len(matches) == 2:
            self.okr_link = matches[0]
            self.okr_id = matches[1]
        else:
            self.okr_link = ""
            self.okr_id = ""
                    
def flip_bool_attribute(obj, attribute):
    original_value = getattr(obj, attribute)
    setattr(obj, attribute, not original_value)
    setattr(obj, attribute, original_value)

def add_run_with_text(paragraph, text, bold=False, font_size=14):
    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    return run
        
def add_paragraph_with_text(text_frame, text, bold=False, font_size=14, level=0, font_color=-1):
    p = text_frame.add_paragraph()
    run = add_run_with_text(p, text, bold, font_size)    
    if font_color != -1:
        run.font.fill.solid()
        run.font.fill.fore_color.rgb = font_color
    p.level = level
    return p

def add_text_block_to_slide(text_frame, text_block_json):
    text_block = json.loads(text_block_json)
    p = None
    for element in text_block['elements']:
        if element.get('is_run', False):
            if p is None:
                raise ValueError("The first element cannot be a run. Add a paragraph first.")
            run = add_run_with_text(p, element['text'], element.get('bold', False), element.get('font_size', 14))
            if 'font_color' in element:
                run.font.fill.solid()
                run.font.fill.fore_color.rgb = RGBColor(*element['font_color'])
        else:
            p = add_paragraph_with_text(text_frame, element['text'], element.get('bold', False), element.get('font_size', 14), element.get('level', 0))
            if 'font_color' in element:
                p.font.fill.solid()
                p.font.fill.fore_color.rgb = RGBColor(*element['font_color'])

def get_goal_by_id(okr_id):
    global goals_dict
    try:
        return goals_dict[okr_id]    
    except KeyError:
        return None
              
def get_theme_goal_by_id(goals):
    for goal in goals:
        if goal.tag == THEME_TAG:
            return goal
    return None

def get_parent_goals_from_alignment(goal):    
    parent_goals = []
    pattern = r"\(weight: \d+(\.\d+)?%, Id: (\d+)\)"
    matches = re.findall(pattern, goal.alignment)
    for match in matches:
        parent_goal = get_goal_by_id(match[1])
        if parent_goal:
            parent_goals.append(parent_goal)
    return parent_goals

# Custom sorting function to ensure goals are showing following this order: 
# Theme, [Outcome,] Objective, [Outcome,] Action [, Theme, [Outcome,] Objective, [Outcome,] Action]
# Notice if there's Objective and Outcome linked to the same Theme, Outcome is shown first.
# If there's Outcome and Action linked to the same Objective, Outcome is shown first.
# The Tuple returned by this function has the follow form: (first level, first sort modifier, second level, second sort modifier, third level)
def goal_sort_key(goal):
    FIRST_PRIORITY = 0
    SECOND_PRIORITY = 1    
    parent_goals = get_parent_goals_from_alignment(goal)
    if goal.object_type == OBJECTIVE_TYPE:
        theme = get_theme_goal_by_id(parent_goals)
        if theme:
            return (theme.row_number, SECOND_PRIORITY, goal.row_number) + (FIRST_PRIORITY,) * 2
    elif goal.object_type == OUTCOME_TYPE:
        theme = get_theme_goal_by_id(parent_goals)
        if theme:
            return (theme.row_number, FIRST_PRIORITY, goal.row_number) + (FIRST_PRIORITY,) * 2
        if parent_goals and len(parent_goals) > 1:
            raise ValueError("More than one parent goal found in alignment for outcome: " + goal.title)
        key = goal_sort_key(parent_goals[0])        
        return (*key[:3], FIRST_PRIORITY, goal.row_number)
    elif goal.object_type == ACTION_TYPE:
        if parent_goals is None or len(parent_goals) == 0:
            raise ValueError("No parent goal found in alignment for action: " + goal.title)
        if parent_goals and len(parent_goals) > 1:
            raise ValueError("More than one parent goal found in alignment for action: " + goal.title)
        key = goal_sort_key(parent_goals[0]) 
        return (*key[:3], SECOND_PRIORITY, goal.row_number)   
    return (goal.row_number,) + (FIRST_PRIORITY,) * 4  # For root-level Themes the code will get to this point
    
def main(source_workbook=SOURCE_WORKBOOK, template_powerpoint=TEMPLATE_POWERPOINT, target_bizplan_powerpoint=TARGET_BIZPLAN_POWERPOINT,
         theme_slide_master=THEME_SLIDE_MASTER, theme_slide_master_layout=THEME_SLIDE_MASTER_LAYOUT,
         okr_slide_master=OKR_SLIDE_MASTER, okr_slide_master_layout=OKR_SLIDE_MASTER_LAYOUT):
    global goals_dict
    wb = load_workbook(source_workbook)
    ws = wb.active

    prs = Presentation(template_powerpoint)

    headers = [cell.value for cell in next(ws.iter_rows(min_row=1, max_row=1))]
    goals_dict = {OKRId(row[headers.index('Id')]).okr_id: VivaGoal(row, headers, idx) for idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True))}
    goals = [VivaGoal(row, headers, idx) for idx, row in enumerate(ws.iter_rows(min_row=2, values_only=True))]
    goals.sort(key=goal_sort_key)

    for goal in goals:    
        if goal.tag == THEME_TAG:
            slide_layout = prs.slide_masters[theme_slide_master].slide_layouts[theme_slide_master_layout]  # Choosing a layout for theme that has only a title object
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = goal.title         
            continue    

        slide_layout = prs.slide_masters[okr_slide_master].slide_layouts[okr_slide_master_layout]  # Choosing a layout for content that has only a title object
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = goal.title
        
        # Add textbox for the OKR attributes
        dimensions = SquareDimensions(left=0.5, top=0.8, width=12, height=3.5)    
        text_box = slide.shapes.add_textbox(dimensions.left, dimensions.top, dimensions.width, dimensions.height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        okr_id = OKRId(goal.okr_id) 

        # Text code with the core attributes of the OKR object
        text_block_json = json.dumps({
            "elements": [
                {"text": "Type: ", "bold": True, "font_size": 18, "level": 1},
                {"text": goal.object_type, "font_size": 18, "level": 1, "is_run": True},
                {"text": ", Metric: ", "bold": True, "font_size": 18, "level": 1, "is_run": True},
                {"text": goal.metric_name, "font_size": 18, "level": 1, "is_run": True},
                {"text": ", Target: ", "bold": True, "font_size": 18, "level": 1, "is_run": True},
                {"text": goal.target, "font_size": 18, "level": 1, "is_run": True},
                {"text": "Owner: ", "bold": True, "font_size": 18, "level": 1},
                {"text": goal.owner, "font_size": 18, "level": 1, "is_run": True},
                {"text": "Schedule: ", "bold": True, "font_size": 18, "level": 1},
                {"text": goal.schedule, "font_size": 18, "level": 1, "is_run": True},
                {"text": "Status: ", "bold": True, "font_size": 18, "level": 1},
                {"text": goal.status, "font_size": 18, "level": 1, "is_run": True}
            ]
        })

        # Add the text block to the slide
        add_text_block_to_slide(text_frame, text_block_json)
        
        # Regular expression to match and remove weight and ID of alignment objective/action
        pattern = r"\(weight: \d+(\.\d+)?%, Id: \d+\)"
        cleaned_alignment = re.sub(pattern, "", goal.alignment)
        
        if goal.object_type == OBJECTIVE_TYPE:        
            image_path = OBJECTIVE_IMAGE
                    
            parts = cleaned_alignment.split(" / ")
            alignment = ""
            mwb = ""

            for part in parts:
                if part.startswith("MWB:"):
                    mwb = part
                else:
                    alignment = part
            if alignment != "":
                p = add_paragraph_with_text(text_frame, "Parent plan theme: ", True, 18, 1)
                add_run_with_text(p, alignment, False, 18)
            if mwb != "":
                add_paragraph_with_text(text_frame, "")
                p = add_paragraph_with_text(text_frame, "Parent MWB alignment: ", True, 18, 1, RGBColor(0, 176, 240))               
                add_run_with_text(p, mwb, False, 18)
                
            # Define the position and size of the title rectangle
            dimensions = SquareDimensions(left=0.5, top=0.3, width=12.5, height=0.75)        

            # Add a dark blue rectangle behind the title when showing Objectives
            title_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, dimensions.left, dimensions.top, dimensions.width, dimensions.height)
            title_rect.fill.solid()
            title_rect.fill.fore_color.rgb = RGBColor(0, 43, 72)  # Blue color
            title_rect.line.color.rgb = RGBColor(0, 0, 255)  # No border
            # Send the rectangle to the back of the title object (title came with the template)
            spTree = slide.shapes._spTree
            spTree.remove(title_rect._element)
            spTree.insert(2, title_rect._element)  # Insert at the beginning to send to back
        else:        
            p = add_paragraph_with_text(text_frame, "Parent objective: ", True, 18, 1)
            add_run_with_text(p, cleaned_alignment, False, 18)     
            if goal.object_type == ACTION_TYPE:
                image_path = INITIATIVE_IMAGE
            else:
                image_path = OUTCOME_IMAGE
        
        # Insert image representing the type of OKR object    
        dimensions = SquareDimensions(left=0.34, top=1.13, width=0.5, height=0.5)    
        pic = slide.shapes.add_picture(image_path, dimensions.left, dimensions.top, dimensions.width, dimensions.height)
        pic.click_action.hyperlink.address = okr_id.okr_link

        line = LineDimensions(left=0.5, top=4, width=12)    
        line_shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, line.left, line.top, line.left + line.width, line.top)

        # Add a text box for the description
        dimensions = SquareDimensions(left=0.5, top=4, width=12, height=3.4)   
        text_box = slide.shapes.add_textbox(dimensions.left, dimensions.top, dimensions.width, dimensions.height)
        text_frame = text_box.text_frame
        add_paragraph_with_text(text_frame, "Description:", bold=True, font_size=18)
        p = add_paragraph_with_text(text_frame, "")
        add_run_with_text(p, goal.description, font_size=14)
        
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.word_wrap = True
        flip_bool_attribute(text_frame, 'word_wrap')
        
    # Save the presentation
    prs.save(target_bizplan_powerpoint)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Transform Viva Goals Excel export into a PowerPoint file.')
    parser.add_argument('--source_workbook', type=str, default='VivaGoals.xlsx', help='Path to the source Excel workbook.')
    parser.add_argument('--template_powerpoint', type=str, default='template.pptx', help='Path to the template PowerPoint file.')
    parser.add_argument('--target_bizplan_powerpoint', type=str, default='bizplan.pptx', help='Path to the target PowerPoint file.')
    parser.add_argument('--theme_slide_master', type=int, default=THEME_SLIDE_MASTER, help='Index of the theme slide master.')
    parser.add_argument('--theme_slide_master_layout', type=int, default=THEME_SLIDE_MASTER_LAYOUT, help='Index of the theme slide master layout.')
    parser.add_argument('--okr_slide_master', type=int, default=OKR_SLIDE_MASTER, help='Index of the OKR slide master.')
    parser.add_argument('--okr_slide_master_layout', type=int, default=OKR_SLIDE_MASTER_LAYOUT, help='Index of the OKR slide master layout.')

    args = parser.parse_args()
    main(source_workbook=args.source_workbook, template_powerpoint=args.template_powerpoint, target_bizplan_powerpoint=args.target_bizplan_powerpoint,
         theme_slide_master=args.theme_slide_master, theme_slide_master_layout=args.theme_slide_master_layout,
         okr_slide_master=args.okr_slide_master, okr_slide_master_layout=args.okr_slide_master_layout)